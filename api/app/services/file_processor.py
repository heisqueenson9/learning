import os
import logging

MAX_PDF_PAGES = 12
MAX_PPTX_SLIDES = 15
MAX_DOCX_PARAGRAPHS = 80
MAX_TEXT_CHARS = 6000  # Cap final output to keep AI fast

def extract_text(file_path: str) -> str:
    _, ext = os.path.splitext(file_path)
    text = ""
    logging.info(f"Extracting text from: {file_path} (ext={ext.lower()})")

    try:
        if ext.lower() == '.pdf':
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            # Limit pages for speed
            pages = reader.pages[:MAX_PDF_PAGES]
            for page in pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
                    if len(text) >= MAX_TEXT_CHARS:
                        break

        elif ext.lower() in ['.docx', '.doc']:
            from docx import Document
            doc = Document(file_path)
            for i, para in enumerate(doc.paragraphs):
                if i >= MAX_DOCX_PARAGRAPHS:
                    break
                text += para.text + "\n"
                if len(text) >= MAX_TEXT_CHARS:
                    break

        elif ext.lower() in ['.pptx', '.ppt']:
            from pptx import Presentation
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                if i >= MAX_PPTX_SLIDES:
                    break
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if len(text) >= MAX_TEXT_CHARS:
                    break

        elif ext.lower() == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read(MAX_TEXT_CHARS)  # Read only up to cap

        else:
            raise ValueError(f"Unsupported file format: {ext}")

    except Exception as e:
        logging.error(f"Error reading file: {e}")
        return ""

    # Trim to max chars and strip whitespace
    return text[:MAX_TEXT_CHARS].strip()
