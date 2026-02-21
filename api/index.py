"""
APEX EduAI Vault — Vercel Serverless API
=========================================
Single-file FastAPI app. All logic is self-contained to avoid
Python import-chain issues on Vercel serverless functions.
"""

import os
import json
import base64
import asyncio
import urllib.parse
import logging
from datetime import datetime
from contextlib import asynccontextmanager
from typing import Optional

from fastapi import FastAPI, Request, UploadFile, Form, Header, HTTPException, Depends, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from starlette.middleware.base import BaseHTTPMiddleware
from pydantic import BaseModel
from sqlalchemy import create_engine, Column, Integer, String, Text, DateTime
from sqlalchemy.orm import sessionmaker, declarative_base, Session

try:
    from dotenv import load_dotenv
    load_dotenv()
except Exception:
    pass

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# DATABASE
# ─────────────────────────────────────────────────────────────────────────────
DATABASE_URL = os.getenv("DATABASE_URL", "sqlite:///./apex.db")
if DATABASE_URL.startswith("postgres://"):
    DATABASE_URL = DATABASE_URL.replace("postgres://", "postgresql://", 1)

try:
    if "sqlite" in DATABASE_URL:
        engine = create_engine(
            DATABASE_URL,
            connect_args={"check_same_thread": False},
            pool_pre_ping=True
        )
    else:
        engine = create_engine(DATABASE_URL, pool_pre_ping=True)
    SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
except Exception as e:
    print(f"[DB] Engine creation failed: {e}")
    engine = None
    SessionLocal = None

Base = declarative_base()


class Payment(Base):
    __tablename__ = "payments"
    id = Column(Integer, primary_key=True, index=True)
    full_name = Column(String(200))
    phone = Column(String(50))
    institution = Column(String(200), nullable=True)
    package = Column(String(100), nullable=True)
    screenshot_url = Column(Text)
    created_at = Column(DateTime, default=datetime.utcnow)


class ExamRecord(Base):
    __tablename__ = "exams"
    id = Column(Integer, primary_key=True, index=True)
    title = Column(String(500))
    level = Column(String(100))
    topic = Column(String(500))
    difficulty = Column(String(100))
    exam_type = Column(String(100))
    questions = Column(Text)
    created_at = Column(DateTime, default=datetime.utcnow)


class GameAnswer(Base):
    __tablename__ = "game_answers"
    id = Column(Integer, primary_key=True, index=True)
    full_name = Column(String(200))
    question = Column(Text)
    answer = Column(Text)
    created_at = Column(DateTime, default=datetime.utcnow)


def init_db():
    try:
        if engine:
            Base.metadata.create_all(bind=engine)
            print("[DB] Tables created/verified.")
    except Exception as e:
        print(f"[DB] Table creation warning: {e}")


def get_db():
    if not SessionLocal:
        raise HTTPException(status_code=503, detail="Database not available.")
    db = SessionLocal()
    try:
        yield db
    finally:
        try:
            db.close()
        except Exception:
            pass


# ─────────────────────────────────────────────────────────────────────────────
# AI ENGINE (inline — no sub-package imports)
# ─────────────────────────────────────────────────────────────────────────────
def _generate_mock_questions(topic: str, num_q: int) -> list:
    """Generate topic-aware MCQ bank. Repeats to fill num_q."""
    t = topic.replace('"', '').strip() or "General Knowledge"

    bank = [
        {"q": f"What is the primary function of {t}?",
         "opts": ["Data storage and retrieval", "Process execution and management",
                  "Network resource allocation", "User interface rendering"],
         "ans": "Process execution and management"},
        {"q": f"Which of the following best defines {t}?",
         "opts": ["A high-level software abstraction", "A low-level hardware component",
                  "A physical user-interface element", "A network routing protocol"],
         "ans": "A high-level software abstraction"},
        {"q": f"What is a fundamental characteristic of {t}?",
         "opts": ["It exhibits static, unchanging behaviour", "It demonstrates dynamic adaptability",
                  "It maintains a completely fixed structure", "It has inherently limited scalability"],
         "ans": "It demonstrates dynamic adaptability"},
        {"q": f"In computing, {t} is primarily associated with:",
         "opts": ["Physical memory management", "Logical data organisation",
                  "Hardware interrupt handling", "Physical network communication"],
         "ans": "Logical data organisation"},
        {"q": f"The concept of {t} was developed primarily to address:",
         "opts": ["Hardware component failures", "Increasing software complexity",
                  "Raw data redundancy", "Physical network latency"],
         "ans": "Increasing software complexity"},
        {"q": f"Which statement about {t} is most accurate?",
         "opts": [f"{t} is only used in embedded systems",
                  f"{t} applies across multiple computing domains",
                  f"{t} was deprecated in the 1990s",
                  f"{t} requires specialised hardware"],
         "ans": f"{t} applies across multiple computing domains"},
        {"q": f"What does {t} primarily seek to optimise?",
         "opts": ["Physical disk I/O throughput", "Logical process efficiency and correctness",
                  "Screen rendering pipeline speed", "Physical memory chip access"],
         "ans": "Logical process efficiency and correctness"},
        {"q": f"Which of these is NOT a property of {t}?",
         "opts": ["Modularity", "Scalability", "Physical hardware dependency", "Abstraction"],
         "ans": "Physical hardware dependency"},
        {"q": f"The theoretical model underlying {t} is best described as:",
         "opts": ["A finite automaton", "An abstract computational model",
                  "A physical circuit diagram", "A relational database schema"],
         "ans": "An abstract computational model"},
        {"q": f"Which academic field most directly contributed to the foundations of {t}?",
         "opts": ["Biochemistry", "Discrete mathematics and logic",
                  "Mechanical engineering", "Organic chemistry"],
         "ans": "Discrete mathematics and logic"},
        {"q": f"When was {t} first formally introduced as a concept?",
         "opts": ["Early 1960s", "Early 1970s", "Mid 1980s", "Late 1990s"],
         "ans": "Mid 1980s"},
        {"q": f"Which organisation played the most significant role in standardising {t}?",
         "opts": ["IEEE", "ISO", "ANSI", "W3C"],
         "ans": "IEEE"},
        {"q": f"Which industry sector makes the most extensive use of {t}?",
         "opts": ["Healthcare", "Information technology", "Financial services", "Manufacturing"],
         "ans": "Information technology"},
        {"q": f"In software engineering, {t} is used primarily for:",
         "opts": ["System performance optimisation", "Security policy management",
                  "Logical code organisation", "Direct database querying"],
         "ans": "Logical code organisation"},
        {"q": f"Which domain gains the greatest productivity improvement from {t}?",
         "opts": ["Machine learning and AI", "Static web development",
                  "Basic networking only", "Simple database design"],
         "ans": "Machine learning and AI"},
        {"q": f"What is the most significant advantage of using {t} in enterprise systems?",
         "opts": ["Drastically increased memory consumption", "Significantly improved maintainability",
                  "Consistently slower execution speed", "Inherently weaker security posture"],
         "ans": "Significantly improved maintainability"},
        {"q": f"Which limitation is most commonly associated with {t}?",
         "opts": ["High initial computational overhead", "Poor or missing documentation",
                  "Complete lack of industry adoption", "Trivially simple implementation"],
         "ans": "High initial computational overhead"},
        {"q": f"A major trade-off when selecting {t} over alternatives is:",
         "opts": ["Processing speed vs. implementation simplicity",
                  "Licensing cost vs. system reliability",
                  "Memory usage vs. execution speed",
                  "Security strength vs. overall usability"],
         "ans": "Memory usage vs. execution speed"},
        {"q": f"What is the core architectural component of a {t} system?",
         "opts": ["The operating system kernel", "The user-facing interface layer",
                  "The core functional module", "The hardware controller"],
         "ans": "The core functional module"},
        {"q": f"{t} is architecturally built upon which fundamental concept?",
         "opts": ["Recursive programming", "Iterative execution",
                  "Abstraction and encapsulation", "Runtime polymorphism"],
         "ans": "Abstraction and encapsulation"},
        {"q": f"How does {t} fundamentally differ from its closest alternative?",
         "opts": ["It consistently performs slower",
                  "It provides a significantly better level of abstraction",
                  "It invariably uses more memory", "It has substantially fewer features"],
         "ans": "It provides a significantly better level of abstraction"},
        {"q": f"What most clearly distinguishes {t} from older legacy methodologies?",
         "opts": ["It is tightly hardware-dependent",
                  "It properly supports modern scalability requirements",
                  "It requires fully manual memory management",
                  "It is fundamentally single-threaded"],
         "ans": "It properly supports modern scalability requirements"},
        {"q": f"In the context of concurrent programming, {t} is primarily associated with:",
         "opts": ["Advanced deadlock prevention mechanisms", "Intentional race condition creation",
                  "Strictly single-threaded execution only",
                  "Exclusively synchronous sequential execution"],
         "ans": "Advanced deadlock prevention mechanisms"},
        {"q": f"When optimising a {t} implementation, which metric is most critical?",
         "opts": ["Total lines of source code", "Asymptotic time and space complexity",
                  "Visual colour scheme of the UI", "Total compressed file size"],
         "ans": "Asymptotic time and space complexity"},
        {"q": f"The scalability of a {t} system is most accurately measured by:",
         "opts": ["The total number of concurrent users",
                  "System throughput and end-to-end latency",
                  "Percentage of test code coverage", "The size of the development team"],
         "ans": "System throughput and end-to-end latency"},
        {"q": f"The concept of modularity within {t} architecture specifically means:",
         "opts": ["Source code cannot be effectively reused",
                  "All components are tightly and rigidly coupled",
                  "Components can be independently developed and maintained",
                  "The entire system forms one indivisible monolith"],
         "ans": "Components can be independently developed and maintained"},
        {"q": f"When comprehensively testing {t}, which approach provides the broadest coverage?",
         "opts": ["Unit testing in complete isolation",
                  "Integration testing of combined components",
                  "Full end-to-end system testing",
                  "A combination of all testing approaches"],
         "ans": "A combination of all testing approaches"},
        {"q": f"A student implementing {t} for the first time should begin by:",
         "opts": ["Immediately writing production code",
                  "Thoroughly reading the official documentation first",
                  "Skipping documentation and experimenting",
                  "Understanding core concepts, reading docs, then writing tests"],
         "ans": "Understanding core concepts, reading docs, then writing tests"},
        {"q": f"When a production {t} system fails, the first diagnostic step should be to:",
         "opts": ["Immediately restart all physical hardware",
                  "Carefully examine the system and application logs",
                  "Perform a complete software reinstallation",
                  "Escalate immediately to the vendor without investigation"],
         "ans": "Carefully examine the system and application logs"},
        {"q": f"When evaluating {t} for a new project, the team should carefully consider:",
         "opts": ["Only the physical office location of the team",
                  "Team expertise, project requirements, and scalability needs",
                  "Only available font and colour choices for the UI",
                  "Only the team's preferred colour palette"],
         "ans": "Team expertise, project requirements, and scalability needs"},
        {"q": f"Which software design pattern most closely aligns with the principles of {t}?",
         "opts": ["Observer pattern", "Abstract Factory pattern",
                  "Singleton pattern", "Strategy pattern"],
         "ans": "Strategy pattern"},
        {"q": f"Which data structure is most commonly used internally by {t} implementations?",
         "opts": ["Simple flat arrays", "Singly linked lists",
                  "Hash tables or maps", "Simple binary trees"],
         "ans": "Hash tables or maps"},
        {"q": f"The architecture of a well-designed {t} system typically includes:",
         "opts": ["Only input/output handlers", "Multiple organised layered components",
                  "A single tightly-coupled monolithic block",
                  "Only external third-party API calls"],
         "ans": "Multiple organised layered components"},
        {"q": f"Which class of problems is {t} best suited to solve?",
         "opts": ["Race conditions in parallel code", "Memory leak detection",
                  "Data inconsistency issues", "Physical network failures"],
         "ans": "Data inconsistency issues"},
        {"q": f"In academic computing curricula, {t} is primarily taught within which course?",
         "opts": ["Systems programming", "Theory of computation",
                  "Computer architecture", "Algorithm design and analysis"],
         "ans": "Algorithm design and analysis"},
        {"q": f"Which of the following best represents a real-world application of {t}?",
         "opts": ["Only image compression", "Only search algorithms",
                  "Only encryption systems", "All of the above and more"],
         "ans": "All of the above and more"},
        {"q": f"What is the relationship between {t} and system performance?",
         "opts": ["It always decreases performance significantly",
                  "It optimises performance through better resource management",
                  "It has no measurable impact on performance",
                  "It only affects network performance"],
         "ans": "It optimises performance through better resource management"},
        {"q": f"Which programming paradigm is most closely related to {t}?",
         "opts": ["Purely functional programming", "Object-oriented programming",
                  "Assembly-level programming", "Scripting-only approaches"],
         "ans": "Object-oriented programming"},
        {"q": f"How is {t} typically validated in a production environment?",
         "opts": ["By visual inspection of source code only",
                  "Through automated testing, monitoring, and logging",
                  "By running the system once without errors",
                  "Only through manual user acceptance testing"],
         "ans": "Through automated testing, monitoring, and logging"},
        {"q": f"What role does documentation play in a well-maintained {t} system?",
         "opts": ["Documentation is optional and rarely useful",
                  "It is essential for maintainability and onboarding",
                  "Documentation only matters for open-source projects",
                  "It slows down development with no real benefit"],
         "ans": "It is essential for maintainability and onboarding"},
        {"q": f"Which characteristic makes {t} particularly preferable in enterprise environments?",
         "opts": ["Its extremely high reliability metrics", "Its typically very low cost",
                  "Its simple initial setup process", "Its minimal external dependencies"],
         "ans": "Its extremely high reliability metrics"},
        {"q": f"The ease of parallelisation in {t} is primarily due to:",
         "opts": ["Its low-level hardware access", "Its clean separation of concerns and modularity",
                  "Its use of a single global state", "Its reliance on synchronous I/O"],
         "ans": "Its clean separation of concerns and modularity"},
        {"q": f"What is the best strategy for scaling a {t}-based system?",
         "opts": ["Vertically scale the database only",
                  "Horizontal scaling with load balancing",
                  "Reduce the feature set to lower load",
                  "Use a single powerful machine indefinitely"],
         "ans": "Horizontal scaling with load balancing"},
        {"q": f"Which principle is most fundamental to the design of {t}?",
         "opts": ["Tight coupling between components", "Single responsibility principle",
                  "Global mutable state", "Avoidance of all abstraction"],
         "ans": "Single responsibility principle"},
        {"q": f"How should errors be handled in a robust {t} implementation?",
         "opts": ["Silently ignored to improve performance",
                  "Gracefully — with logging, fallbacks, and user-friendly messages",
                  "By crashing the entire system immediately",
                  "Only caught at the outermost layer of the application"],
         "ans": "Gracefully — with logging, fallbacks, and user-friendly messages"},
        {"q": f"What is the primary security consideration when implementing {t}?",
         "opts": ["Making the user interface as simple as possible",
                  "Validation, authentication, and least-privilege access",
                  "Using only free and open-source libraries",
                  "Avoiding any external API calls"],
         "ans": "Validation, authentication, and least-privilege access"},
        {"q": f"Which metric best reflects the code quality of a {t} implementation?",
         "opts": ["Number of lines of code", "Test coverage, code review, and maintainability index",
                  "The number of features implemented", "The speed of initial development"],
         "ans": "Test coverage, code review, and maintainability index"},
        {"q": f"How does {t} handle backward compatibility?",
         "opts": ["It ignores backward compatibility entirely",
                  "It maintains stability through versioning and deprecation strategies",
                  "By removing old features immediately without notice",
                  "It never changes once deployed"],
         "ans": "It maintains stability through versioning and deprecation strategies"},
        {"q": f"What is the recommended approach for deploying {t} to a cloud environment?",
         "opts": ["Manual FTP upload of files",
                  "CI/CD pipelines with containerisation and infrastructure-as-code",
                  "Only deploying during low-traffic periods manually",
                  "Emailing zip files to the server administrator"],
         "ans": "CI/CD pipelines with containerisation and infrastructure-as-code"},
        {"q": f"Which of the following best describes the impact of {t} on developer productivity?",
         "opts": ["It consistently reduces developer productivity",
                  "It significantly improves productivity through reusability and clarity",
                  "It has no measurable impact on productivity",
                  "Only senior developers benefit from it"],
         "ans": "It significantly improves productivity through reusability and clarity"},
    ]

    questions = []
    for i in range(num_q):
        item = bank[i % len(bank)]
        questions.append({
            "id": i + 1,
            "question": item["q"],
            "options": item["opts"],
            "answer": item["ans"],
        })
    return questions


def _try_pollinations_ai(topic: str, difficulty: str, num_q: int, context: str = "") -> Optional[list]:
    """Try to get questions from Pollinations AI (free, no key)."""
    try:
        import httpx
        prompt = (
            f"Generate exactly {num_q} multiple choice questions about '{topic}'. "
            f"Difficulty: {difficulty}. "
            f"Return ONLY valid JSON with this exact format: "
            f'{{\"questions\": [{{\"id\": 1, \"question\": \"...\", \"options\": [\"A\", \"B\", \"C\", \"D\"], \"answer\": \"A\"}}]}}. '
            f"No markdown, no extra text. Context: {context[:1000]}"
        )
        url = f"https://text.pollinations.ai/{urllib.parse.quote(prompt)}?json=true"
        with httpx.Client(timeout=55.0) as client:
            resp = client.get(url)
            resp.raise_for_status()
            content = resp.text.replace("```json", "").replace("```", "").strip()
            start = content.find("{")
            end = content.rfind("}") + 1
            if start != -1 and end > start:
                parsed = json.loads(content[start:end])
                qs = parsed.get("questions", [])
                if len(qs) >= max(5, num_q // 5):
                    # Add sequential IDs
                    for i, q in enumerate(qs):
                        q["id"] = i + 1
                    return qs
    except Exception as e:
        print(f"[AI] Pollinations failed: {e}")
    return None


def generate_questions(topic: str, difficulty: str, num_q: int, text_context: str = "") -> dict:
    """Generate a quiz bank. Tries free AI first, falls back to mock."""
    # Try free AI
    qs = _try_pollinations_ai(topic, difficulty, num_q, text_context)

    if qs:
        # Pad to num_q if needed
        if len(qs) < num_q:
            mock_extra = _generate_mock_questions(topic, num_q - len(qs))
            for i, q in enumerate(mock_extra):
                q["id"] = len(qs) + i + 1
            qs.extend(mock_extra)
        qs = qs[:num_q]
    else:
        qs = _generate_mock_questions(topic, num_q)

    return {
        "title": f"{topic} — {num_q}-Question Assessment",
        "questions": qs,
    }


# ─────────────────────────────────────────────────────────────────────────────
# FILE TEXT EXTRACTION (inline)
# ─────────────────────────────────────────────────────────────────────────────
def extract_text_from_file(file_path: str) -> str:
    """Extract text from uploaded file. Returns empty string on failure."""
    MAX_CHARS = 4000
    _, ext = os.path.splitext(file_path.lower())
    text = ""
    try:
        if ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            for page in reader.pages[:10]:
                t = page.extract_text() or ""
                text += t + "\n"
                if len(text) >= MAX_CHARS:
                    break
        elif ext in (".docx", ".doc"):
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs[:80]:
                text += para.text + "\n"
                if len(text) >= MAX_CHARS:
                    break
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides[:15]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if len(text) >= MAX_CHARS:
                    break
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read(MAX_CHARS)
    except Exception as e:
        print(f"[Extract] Error: {e}")
    return text[:MAX_CHARS].strip()


# ─────────────────────────────────────────────────────────────────────────────
# FASTAPI APP
# ─────────────────────────────────────────────────────────────────────────────
ADMIN_SECRET = os.getenv("ADMIN_SECRET", "FlameFlame@99")


@asynccontextmanager
async def lifespan(application: FastAPI):
    init_db()
    yield


app = FastAPI(title="APEX EduAI API", lifespan=lifespan)


class SecurityHeaders(BaseHTTPMiddleware):
    async def dispatch(self, request: Request, call_next):
        response = await call_next(request)
        response.headers["X-Content-Type-Options"] = "nosniff"
        response.headers["X-Frame-Options"] = "DENY"
        return response


app.add_middleware(SecurityHeaders)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.exception_handler(Exception)
async def global_handler(request: Request, exc: Exception):
    print(f"[ERROR] {type(exc).__name__}: {exc}")
    return JSONResponse(
        status_code=500,
        content={"detail": "Internal server error", "error": str(exc)},
    )


# ── Health / Root ──────────────────────────────────────────────────────────
@app.get("/")
def root():
    return {"message": "APEX EduAI API — Active ✓"}


@app.get("/health")
def health():
    return {"status": "ok"}


@app.get("/api/v1/ai-status")
def ai_status():
    return {"mode": "free", "ready": True, "failed": False}


# ── Upload Payment ─────────────────────────────────────────────────────────
@app.post("/api/v1/upload-payment")
async def upload_payment(
    full_name: str = Form(...),
    phone: str = Form(...),
    institution: str = Form(None),
    package: str = Form(None),
    duration: str = Form(None),
    screenshot: UploadFile = File(...),
    db: Session = Depends(get_db),
):
    try:
        raw = await screenshot.read()
        encoded = base64.b64encode(raw).decode("utf-8")
        mime = screenshot.content_type or "image/jpeg"
        data_uri = f"data:{mime};base64,{encoded}"
        record = Payment(
            full_name=full_name.strip(),
            phone=phone.strip(),
            institution=(institution or "Not Specified").strip(),
            package=(package or "Unknown").strip(),
            screenshot_url=data_uri,
        )
        db.add(record)
        db.commit()
    except Exception as e:
        print(f"[Payment] Save error (granting access anyway): {e}")
        try:
            db.rollback()
        except Exception:
            pass
    # Always grant access — any screenshot upload grants access
    return {"access": True, "message": "Access Granted"}


# ── Admin ──────────────────────────────────────────────────────────────────
def _admin_check(x_admin_key: str = Header(default="")):
    if x_admin_key != ADMIN_SECRET:
        raise HTTPException(status_code=403, detail="Invalid admin key.")


class AdminLoginIn(BaseModel):
    phone: str
    password: str


@app.post("/api/v1/admin/login")
def admin_login(payload: AdminLoginIn):
    if payload.password == ADMIN_SECRET:
        return {"message": "Admin login successful", "token": ADMIN_SECRET}
    raise HTTPException(status_code=401, detail="Invalid credentials.")


@app.get("/api/v1/admin/payments")
def admin_payments(db: Session = Depends(get_db), _=Depends(_admin_check)):
    try:
        rows = db.query(Payment).order_by(Payment.created_at.desc()).all()
        return [
            {
                "id": r.id,
                "full_name": r.full_name,
                "phone": r.phone,
                "institution": r.institution,
                "package": r.package,
                "screenshot_url": r.screenshot_url,
                "created_at": r.created_at.isoformat() if r.created_at else None,
            }
            for r in rows
        ]
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/api/v1/admin/payments/{payment_id}")
def delete_payment(payment_id: int, db: Session = Depends(get_db), _=Depends(_admin_check)):
    p = db.query(Payment).filter(Payment.id == payment_id).first()
    if p:
        db.delete(p)
        db.commit()
    return {"message": "Deleted"}


@app.get("/api/v1/admin/games")
def admin_games(db: Session = Depends(get_db)):
    try:
        rows = db.query(GameAnswer).order_by(GameAnswer.created_at.desc()).all()
        return [
            {
                "id": r.id,
                "full_name": r.full_name,
                "question": r.question,
                "answer": r.answer,
                "created_at": r.created_at.isoformat() if r.created_at else None,
            }
            for r in rows
        ]
    except Exception as e:
        return []


# ── Exam: Generate ─────────────────────────────────────────────────────────
@app.post("/api/v1/exams/generate")
async def generate_exam(
    level: str = Form(...),
    exam_type: str = Form(...),
    difficulty: str = Form(...),
    topic: str = Form(None),
    file: Optional[UploadFile] = File(None),
    db: Session = Depends(get_db),
):
    topic_name = (topic or "").strip()
    text_context = ""

    # Handle file upload
    if file and file.filename:
        try:
            raw = await file.read()
            tmp = f"/tmp/apex_{int(datetime.utcnow().timestamp())}_{file.filename}"
            with open(tmp, "wb") as fout:
                fout.write(raw)
            text_context = await asyncio.to_thread(extract_text_from_file, tmp)
            if not topic_name:
                topic_name = file.filename
            try:
                os.remove(tmp)
            except Exception:
                pass
        except Exception as e:
            print(f"[Generate] File error: {e}")
            if not topic_name:
                topic_name = file.filename or "Uploaded Document"

    if not topic_name:
        raise HTTPException(status_code=400, detail="A topic or file is required.")

    if not text_context:
        text_context = (
            f"Topic: {topic_name}. Level: {level}. "
            f"Generate 100 comprehensive multiple-choice questions."
        )

    print(f"[Generate] Generating 100 Qs for: {topic_name}")
    result = await asyncio.to_thread(generate_questions, topic_name, difficulty, 100, text_context)

    # Persist to DB (best-effort)
    exam_id = 0
    exam_created = datetime.utcnow().isoformat()
    try:
        rec = ExamRecord(
            title=f"{topic_name} — {exam_type}",
            level=level,
            topic=topic_name,
            difficulty=difficulty,
            exam_type=exam_type,
            questions=json.dumps(result),
        )
        db.add(rec)
        db.commit()
        db.refresh(rec)
        exam_id = rec.id
        exam_created = rec.created_at.isoformat() if rec.created_at else exam_created
    except Exception as e:
        print(f"[Generate] DB save warn: {e}")
        try:
            db.rollback()
        except Exception:
            pass

    return {
        "id": exam_id,
        "title": f"{topic_name} — {exam_type}",
        "questions": result,   # Already a dict — frontend receives it directly
        "created_at": exam_created,
    }


# ── Exam: History ──────────────────────────────────────────────────────────
@app.get("/api/v1/exams/history")
def exam_history(db: Session = Depends(get_db)):
    try:
        rows = db.query(ExamRecord).order_by(ExamRecord.created_at.desc()).all()
        result = []
        for r in rows:
            try:
                cat = r.created_at.isoformat() if r.created_at else None
            except Exception:
                cat = None
            result.append({
                "id": r.id,
                "title": r.title or "Untitled",
                "topic": r.topic or "",
                "level": r.level or "",
                "exam_type": r.exam_type or "",
                "difficulty": r.difficulty or "",
                "created_at": cat,
            })
        return {"exams": result, "total": len(result)}
    except Exception as e:
        print(f"[History] Error: {e}")
        return {"exams": [], "total": 0}


# ── Games ──────────────────────────────────────────────────────────────────
class GameAnswerIn(BaseModel):
    full_name: str
    question: str
    answer: str


@app.post("/api/v1/games/answer")
def submit_answer(payload: GameAnswerIn, db: Session = Depends(get_db)):
    try:
        ga = GameAnswer(
            full_name=payload.full_name,
            question=payload.question,
            answer=payload.answer,
        )
        db.add(ga)
        db.commit()
    except Exception as e:
        print(f"[Games] Save error: {e}")
    return {"message": "Saved"}


# ── Stubs (keep frontend from 404ing on optional endpoints) ────────────────
@app.post("/api/v1/auth/adult-game-log-v2")
def stub_log():
    return {"status": "success"}


@app.post("/api/v1/auth/upload-avatar")
def stub_avatar():
    return {"avatar_url": ""}
